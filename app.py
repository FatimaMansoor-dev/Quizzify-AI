from flask import (
    Flask,
    send_file,
    render_template,
    request,
    jsonify,
    session,
    redirect,
    url_for,
)
import requests
import json
from datetime import datetime, timedelta
import random
from flask_mail import Mail, Message
from groq import Groq
from flask import Flask, request, jsonify
import fitz
import google.generativeai as genai

from pptx import Presentation
from pptx.util import Inches, Pt
import logging
import re

# from youtube_transcript_api import YouTubeTranscriptApi
import speech_recognition as sr
from pydub import AudioSegment
import io
from bs4 import BeautifulSoup
import time

# New imports for Google token verification
import google.oauth2.id_token
import google.auth.transport.requests

# Import for our custom decorator
from functools import wraps

# ===== New Imports for TTS (using gTTS and playsound) =====
from gtts import gTTS
import tempfile
import os
import playsound

# Load environment variables from .env file
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__)

# Set the secret key from the environment variable
app.secret_key = os.environ.get("SECRET_KEY", "default-secret-key")

app.config["GOOGLE_CLIENT_ID"] = os.environ.get("GOOGLE_CLIENT_ID")


@app.context_processor
def inject_config():
    return dict(config=app.config)


# ---------------------
# Alternative TTS Function using gTTS and playsound
# ---------------------
def speak(text):
    try:
        tts = gTTS(text=text, lang="en")
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as fp:
            temp_filename = fp.name
            tts.write_to_fp(fp)
        playsound.playsound(temp_filename)
        os.remove(temp_filename)
    except Exception as e:
        print("TTS error:", e)


# ---------------------
# CORS, Mail, Groq Setup
# ---------------------
from flask_cors import CORS

CORS(app)

# Mail configuration using environment variables
app.config["MAIL_SERVER"] = os.environ.get("MAIL_SERVER", "smtp.gmail.com")
app.config["MAIL_PORT"] = int(os.environ.get("MAIL_PORT", 587))
app.config["MAIL_USE_TLS"] = os.environ.get("MAIL_USE_TLS", "True").lower() in [
    "true",
    "1",
]
app.config["MAIL_USERNAME"] = os.environ.get("MAIL_USERNAME")
app.config["MAIL_PASSWORD"] = os.environ.get("MAIL_PASSWORD")
app.config["MAIL_DEFAULT_SENDER"] = os.environ.get(
    "MAIL_DEFAULT_SENDER", "benzenering032@gmail.com"
)
mail = Mail(app)

# Groq client using API key from environment variables
groq_api_key = os.environ.get("GROQ_API_KEY")
client = Groq(api_key=groq_api_key)


genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
gemini_model = genai.GenerativeModel("gemini-2.0-flash")
# ---------------------
# Firebase / JSON config
# ---------------------
FIREBASE_URL = os.environ.get("FIREBASE_URL")

# Load character mapping from JSON
with open("char_key_mapping.json", "r") as f:
    char_mapping = json.load(f)


# ---------------------
# Custom login_required Decorator
# ---------------------
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if "email" not in session:
            # Redirect to the sign-in page if not logged in
            return redirect(url_for("sign"))
        return f(*args, **kwargs)

    return decorated_function


# ---------------------
# Helper Functions
# ---------------------
def encrypt_password(password):
    return "".join(char_mapping.get(letter, letter) for letter in password)


def decrypt_password(encrypted_password):
    reverse_mapping = {v: k for k, v in char_mapping.items()}
    return "".join(reverse_mapping.get(letter, letter) for letter in encrypted_password)


def get_user_data(email):
    """Fetch user data from Firebase by email."""
    try:
        response = requests.get(FIREBASE_URL)
        # print(response.json())  # Debug: print all users

        if response.status_code != 200:
            raise Exception("Failed to fetch user data from Firebase")

        users = response.json()
        if not users:
            print("im here")
            return None  # No users found

        if isinstance(users, dict):
            print("ahha")
            for user_id, user_info in users.items():
                if user_info.get("email") == email:
                    return user_info
        print("not here")
        return None
    except Exception as e:
        print(f"Error fetching user data: {e}")
        return None


def get_badge_message(quiz_count):
    badges = {
        1: "First Step",
        5: "Quiz Master",
        50: "Quiz Enthusiast",
        100: "Quiz Prodigy",
        500: "Quiz Mastermind",
    }
    earned_badges = []
    # Add badges earned so far
    for threshold in sorted(badges.keys()):
        if quiz_count >= threshold:
            earned_badges.append(f"{badges[threshold]} Badge")
    # Check if user is 1 quiz away from next badge
    for threshold in sorted(badges.keys()):
        if quiz_count + 1 == threshold:
            earned_badges.append(
                f"You are one quiz away from getting the '{badges[threshold]}' badge."
            )
            break
    return (
        earned_badges
        if earned_badges
        else ["Keep going! Your first badge is just one quiz away!"]
    )


def prepare_user_profile(user_data):
    first_name = user_data.get("first_name", "User")
    last_name = user_data.get("last_name", "")
    quiz_attempts = user_data.get("quiz_attempts", [])
    print("quiz attempts:", quiz_attempts)
    # For calendar display, we only need the date strings.
    quiz_dates = [
        attempt.get("date").split(" ")[0]
        for attempt in quiz_attempts
        if "date" in attempt
    ]
    last_score = quiz_attempts[-1].get("score") if quiz_attempts else None
    max_score = max(
        (attempt.get("score", 0) for attempt in quiz_attempts), default=None
    )
    difficulty_counts = {"easy": 0, "medium": 0, "hard": 0}
    for attempt in quiz_attempts:
        difficulty = attempt.get("difficulty", "").lower()
        if difficulty in difficulty_counts:
            difficulty_counts[difficulty] += 1
    quiz_count = len(quiz_attempts)
    message = get_badge_message(quiz_count)

    return {
        "message": message,
        "first_name": first_name,
        "last_name": last_name,
        "quizes": quiz_count,
        "quiz_attempts": quiz_attempts,  # full data for quiz cards
        "quiz_dates": quiz_dates,  # date-only list for calendar & streak tracker
        "last_score": last_score,
        "max_score": max_score,
        "difficulty_counts": difficulty_counts,
    }


import re


def store_quiz_in_firebase(
    email, quiz_content, difficulty, qtype, score=None, extracted_topic=None
):
    """
    Stores the entire quiz attempt in Firebase for the given email.
    Now also stores the extracted topic.
    """
    # 1) Fetch all users from Firebase
    resp = requests.get(FIREBASE_URL)
    if resp.status_code != 200:
        raise Exception("Failed to fetch user data from Firebase")
    users = resp.json() or {}

    # 2) Find the user by email
    user_id = None
    user_info = None
    for uid, info in users.items():
        if info.get("email") == email:
            user_id = uid
            user_info = info
            break
    if not user_info:
        raise Exception("User not found in Firebase")

    # 3) Prepare the new quiz attempt with all required fields
    quiz_attempts = user_info.get("quiz_attempts", [])
    current_date = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    new_attempt = {
        "date": current_date,
        "difficulty": difficulty,
        "score": score,  # Will be None for a newly generated quiz
        "type": qtype,
        "quiz_content": quiz_content,
        "topic": extracted_topic,  # storing the topic along with quiz attempt
    }
    print(extracted_topic)
    quiz_attempts.append(new_attempt)

    # 4) Update the user's quiz_attempts in Firebase
    update_url = FIREBASE_URL.replace(".json", "") + f"/{user_id}.json"
    update_data = {"quiz_attempts": quiz_attempts}
    patch_resp = requests.patch(update_url, json=update_data)
    if patch_resp.status_code != 200:
        raise Exception("Failed to update user data in Firebase")


# ---------------------
# Existing Routes
# ---------------------
@app.route("/")
def home():
    return render_template("home.html", message=None)


@app.route("/sign")
def sign():
    return render_template("index.html", message=None)


@app.route("/login", methods=["POST"])
def login():
    email = request.form.get("email")
    password = request.form.get("password")

    if not email or not password:
        return render_template("index.html", message="Please fill in both fields.")
    if len(password) < 9:
        return render_template(
            "index.html", message="Password must be at least 9 characters long."
        )

    try:
        user_data = get_user_data(email)
        print("got it", user_data)
        if user_data:
            decrypted_password = decrypt_password(user_data["password"])
            print(decrypted_password)
            if password == decrypted_password:
                print("yes")
                session["email"] = email
                profile_data = prepare_user_profile(user_data)
                print("yoo")
                return render_template(
                    "userProfile.html",
                    email=email,
                    user_data=json.dumps(user_data),
                    **profile_data,
                )
        return render_template("index.html", message="Invalid email or password.")
    except Exception as e:
        return render_template("index.html", message=f"An error occurred: {e}")


@app.route("/signup", methods=["POST"])
def signup():
    email = request.form.get("email")
    password = request.form.get("password")
    first_name = request.form.get("first_name")
    last_name = request.form.get("last_name")
    if not email or not password or not first_name or not last_name:
        return render_template("index.html", message="Please fill in all fields.")
    if len(password) < 9:
        return render_template(
            "index.html", message="Password must be at least 9 characters long."
        )
    if get_user_data(email):
        return render_template(
            "index.html", message="Email already exists. Please use a different email."
        )
    encrypted = encrypt_password(password)
    data = {
        "email": email,
        "password": encrypted,
        "first_name": first_name,
        "last_name": last_name,
        "quiz_attempts": [],
    }
    resp = requests.post(FIREBASE_URL, json=data)
    if resp.status_code == 200:
        session["email"] = email
        profile_data = prepare_user_profile(data)
        return render_template(
            "userProfile.html", email=email, user_data=json.dumps(data), **profile_data
        )
    return render_template("index.html", message="Failed to create an account.")


@app.route("/verify-otp", methods=["POST"])
def verify_otp():
    user_otp = request.form.get("otp_combined")
    stored_otp = session.get("otp")
    email = session.get("pending_email")
    encrypted_password = session.get("pending_password")
    first_name = session.get("pending_first_name")
    last_name = session.get("pending_last_name")

    if not user_otp or not stored_otp:
        return render_template(
            "verify_otp.html", email=email, message="OTP is required."
        )

    try:
        if int(user_otp) == stored_otp:
            data = {
                "email": email,
                "password": encrypted_password,
                "first_name": first_name,
                "last_name": last_name,
                "quiz_attempts": [],
            }
            response = requests.post(FIREBASE_URL, json=data)
            if response.status_code == 200:
                session["email"] = email
                session.pop("pending_email", None)
                session.pop("pending_password", None)
                session.pop("pending_first_name", None)
                session.pop("pending_last_name", None)
                profile_data = prepare_user_profile(data)
                return render_template(
                    "userProfile.html",
                    email=email,
                    user_data=json.dumps(data),
                    **profile_data,
                )
            return render_template(
                "verify_otp.html", email=email, message="Failed to create an account."
            )
        else:
            return render_template(
                "verify_otp.html", email=email, message="Invalid OTP. Please try again."
            )
    except Exception as e:
        return render_template(
            "verify_otp.html", email=email, message=f"An error occurred: {e}"
        )


# ---------------------
# Protected Routes (require login)
# ---------------------
@app.route("/userProfile", methods=["GET"])
@login_required
def user_profile():
    email = session.get("email")
    if not email:
        return jsonify({"error": "Email is required"}), 400

    try:
        user_data = get_user_data(email)
        if not user_data:
            return jsonify({"error": "User not found"}), 404

        profile_data = prepare_user_profile(user_data)
        return render_template(
            "userProfile.html",
            email=email,
            user_data=json.dumps(user_data),
            **profile_data,
        )
    except Exception as e:
        return jsonify({"error": f"An error occurred: {e}"}), 500


@app.route("/options")
@login_required
def options():
    email = session.get("email")
    pres = request.args.get("pres")
    print(pres)
    if pres == "ppt":
        return render_template("presentation.html", email=email)
    else:
        return render_template("options.html", email=email)


@app.route("/youtube")
@login_required
def youtube():
    email = session.get("email")
    return render_template("youtube.html", email=email)


@app.route("/pdf")
@login_required
def pdf():
    email = session.get("email")
    return render_template("pdf.html", email=email)


@app.route("/website")
@login_required
def website():
    email = session.get("email")
    return render_template("website.html", email=email)


@app.route("/topic")
@login_required
def topic():
    email = session.get("email")
    return render_template("topic.html", email=email)


@app.route("/voice")
@login_required
def voice():
    email = session.get("email")
    return render_template("voice.html", email=email)


@app.route("/quiz")
@login_required
def quiz():
    message = request.args.get("message", "No message provided")
    email = session.get("email")
    return render_template("quiz.html", message=message, email=email)


@app.route("/blank")
@login_required
def blank():
    message = request.args.get("message", "No message provided")
    email = session.get("email")
    return render_template("blank.html", email=email, message=message)


@app.route("/truefalse")
@login_required
def truefalse():
    message = request.args.get("message", "No message provided")
    email = session.get("email")
    return render_template("truefalse.html", email=email, message=message)


@app.route("/qa")
@login_required
def generalqa():
    message = request.args.get("message", "No message provided")
    email = session.get("email")
    return render_template("qa.html", email=email, message=message)


@app.route("/generateOnTopic", methods=["POST"])
@login_required
def generate_on_topic():
    data = request.get_json()
    email = session.get("email")
    topic = data.get("topic")
    qtype = data.get("type")
    difficulty = data.get("difficulty")
    is_mail = data.get("is_email")

    if qtype == "MCQs":
        prompt = f"""
        Generate exact 10 MCQs quiz with {difficulty} difficulty level on {topic}. Also create a one liner topic for the quiz. If you are generating the quiz from within the data provided to you, then also provide that line as source of answer. If user did not provide you with complete data then do not include this source field.
        Return in this format:
        Topic : [topic]
        **Question 1:** [question]?
        A) [option 1]
        B) [option 2]
        C) [option 3]
        D) [option 4]
        **Answer:** B)
        soure: [should be the exact line that contains the answer and the place or website that you got this answer from.]
        """
    elif qtype in ["blanks", "fillintheblank"]:
        prompt = f"""Generate exact 10 meaningful and logical fill‑in‑the‑blank questions based on {topic} of {difficulty} difficulty level.  
- Each question must be a single sentence with exactly one blank represented by “_____”.  
- After each question, specify **Answer:** with only the phrase that fills the blank.  
- If you were provided the exact sentence containing the answer, include a **Source:** field with that exact sentence; otherwise omit it.

Always reply in this format, for example:

Topic: {topic}

**Question 1:**  _____ is a part of the digestive system is responsible for absorbing most of the nutrients from the food we eat 
**Answer:** small intestine  
soure: [should be the exact line that contains the answer]

… and so on through Question 10."""

    elif qtype in ["truefalse", "true/false"]:
        prompt = f"""
        Generate exact 10 meaningful {difficulty} True/False questions out of which some will be true and some will be false on {topic} + their answers. Also create a one liner topic for the quiz. If you are generating the quiz from within the data provided to you, then also provide that line as source of answer. If user did not provide you with complete data then do not include this source field.
        Format:
        Topic : [topic]
        **Question 1:** [statement]
        **Answer:** [True/False]
        soure: [should be the exact line that contains the answer]
        """
    else:
        prompt = f"""
        Generate exact 10 {difficulty} descriptive Q&A on {topic}. Also create a one liner topic for the quiz.
        Format:
        Topic: [topic]
        **Question 1:** [question]
        **Answer:** [answer]
        soure: [should be the exact line that contains the answer]
        """

    try:
        completion = client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {
                    "role": "system",
                    "content": "You are responsible to generate quizzes.",
                },
                {"role": "user", "content": prompt},
            ],
            temperature=0.3,
            max_tokens=1024,
            top_p=1,
            stream=True,
            stop=None,
        )

        extended_answer = ""
        for response_chunk in completion:
            extended_answer += response_chunk.choices[0].delta.content or ""
        print(extended_answer)

        # Extract the topic from the LM response using a regex search.
        match = re.search(r"Topic\s*:\s*(.*)", extended_answer)
        if match:
            extracted_topic = match.group(1).strip()
        else:
            extracted_topic = topic  # fallback to the original topic if not found

        print(extracted_topic)

        store_quiz_in_firebase(
            email,
            extended_answer,
            difficulty,
            qtype,
            score=None,
            extracted_topic=extracted_topic,
        )

        # If is_mail is True, send the quiz to the user's email
        if is_mail:
            try:
                msg = Message("Your Quiz from QuizifyAI", recipients=[email])
                msg.body = f"Hello,\n\nYour quiz is ready:\n\n{extended_answer}\n\nEnjoy your quiz!"
                mail.send(msg)
            except Exception as e:
                print("Error sending email:", e)

        return jsonify(
            {
                "message": extended_answer,
                "email": email,
                "difficulty": difficulty,
                "type": qtype,
            }
        )
    except Exception as e:
        return (
            jsonify({"error": f"An error occurred while generating the quiz: {e}"}),
            500,
        )


@app.route("/insertResults", methods=["POST", "GET"])
@login_required
def insert_results():
    if request.method == "POST":
        data = request.json
        email = session.get("email")
        score = data.get("score")
        difficulty = data.get("difficulty")

        if not email:
            return jsonify({"error": "Email is required"}), 400

        try:
            response = requests.get(FIREBASE_URL)
            if response.status_code != 200:
                return (
                    jsonify({"error": "Failed to fetch user data from Firebase"}),
                    500,
                )

            users = response.json()
            user_id = None
            user_data = None

            for uid, user in users.items():
                if user.get("email") == email:
                    user_id = uid
                    user_data = user
                    break

            if not user_data:
                return jsonify({"error": "User not found"}), 404

            current_date = datetime.now() + timedelta(hours=2)
            formatted_date = current_date.strftime("%Y-%m-%d %H:%M:%S")
            quiz_attempts = user_data.get("quiz_attempts", [])

            updated = False
            for attempt in reversed(quiz_attempts):
                if attempt.get("score") is None:
                    attempt["score"] = score
                    attempt["date"] = formatted_date
                    attempt["difficulty"] = difficulty
                    updated = True
                    break

            if not updated:
                quiz_attempts.append(
                    {"score": score, "date": formatted_date, "difficulty": difficulty}
                )

            update_url = FIREBASE_URL.replace(".json", "") + f"/{user_id}.json"
            update_data = {"quiz_attempts": quiz_attempts}
            update_response = requests.patch(update_url, json=update_data)

            if update_response.status_code != 200:
                return jsonify({"error": "Failed to update user data in Firebase"}), 500

            return jsonify({"success": "Score saved"}), 200
        except Exception as e:
            print(f"Error storing quiz results: {e}")
            return (
                jsonify({"error": "An error occurred while storing quiz results"}),
                500,
            )


# ---------------------
# Updated /transcribe Endpoint
# ---------------------
@app.route("/transcribe", methods=["POST"])
@login_required
def transcribe():
    try:
        if "file" not in request.files:
            speak("No file uploaded")
            return jsonify({"error": "No file uploaded"}), 400

        audio_file = request.files["file"]
        file_ext = audio_file.filename.rsplit(".", 1)[-1].lower()
        valid_formats = ["mp3", "wav", "m4a", "flac", "ogg", "aac"]

        if file_ext not in valid_formats:
            speak("Unsupported file format")
            return jsonify({"error": "Unsupported file format"}), 400

        audio_bytes = audio_file.read()
        audio_io = io.BytesIO(audio_bytes)
        audio = AudioSegment.from_file(audio_io, format=file_ext)
        wav_io = io.BytesIO()
        audio.export(wav_io, format="wav", parameters=["-acodec", "pcm_s16le"])
        wav_io.seek(0)

        recognizer = sr.Recognizer()
        recognizer.pause_threshold = 1
        recognizer.energy_threshold = 150

        with sr.AudioFile(wav_io) as source:
            audio_data = recognizer.record(source)
            text = recognizer.recognize_google(audio_data)
            speak("The transcribed text is: " + text)

        return jsonify({"text": text})
    except sr.UnknownValueError:
        speak("I did not understand the audio.")
        return jsonify({"error": "Could not understand audio"}), 400
    except sr.RequestError:
        speak("Speech recognition service error.")
        return jsonify({"error": "Speech recognition service error"}), 500
    except Exception as e:
        speak("An error occurred during transcription.")
        return jsonify({"error": str(e)}), 500


# ---------------------
# New /listen Endpoint for Microphone Input
# ---------------------
@app.route("/listen", methods=["GET"])
@login_required
def listen_audio():
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        speak("Listening, please speak now.")
        recognizer.pause_threshold = 1
        recognizer.energy_threshold = 150
        try:
            audio = recognizer.listen(source, timeout=4)
            speak("Understanding...")
            text = recognizer.recognize_google(audio, language="en-in")
            speak("You said: " + text)
            return jsonify({"text": text})
        except sr.UnknownValueError:
            speak("I did not understand that. Could you please repeat?")
            return jsonify({"error": "Could not understand audio"}), 400
        except sr.RequestError:
            speak("Speech recognition service error.")
            return jsonify({"error": "Speech recognition service error"}), 500
        except Exception as e:
            speak("An error occurred. Please try again.")
            return jsonify({"error": str(e)}), 500


# ---------------------
# Google Authentication Route (unprotected)
# ---------------------
@app.route("/auth/google", methods=["POST"])
def auth_google():
    data = request.get_json()
    token = data.get("token")
    if not token:
        return render_template("index.html", message="Token not provided"), 400

    # CLIENT_ID now from env file
    CLIENT_ID = os.environ.get("GOOGLE_CLIENT_ID")
    try:
        request_adapter = google.auth.transport.requests.Request()
        id_info = google.oauth2.id_token.verify_oauth2_token(
            token, request_adapter, CLIENT_ID
        )

        email = id_info.get("email")
        first_name = id_info.get("given_name", "")
        last_name = id_info.get("family_name", "")

        user_data = get_user_data(email)

        if not user_data:
            new_user = {
                "email": email,
                "password": "",
                "first_name": first_name,
                "last_name": last_name,
                "google": True,
                "quiz_attempts": [],
            }
            response = requests.post(FIREBASE_URL, json=new_user)
            if response.status_code == 200:
                user_data = new_user
            else:
                return (
                    render_template(
                        "index.html", message="Failed to create user account"
                    ),
                    500,
                )

        session["email"] = email

        profile_data = prepare_user_profile(user_data)
        return render_template(
            "userProfile.html",
            email=email,
            user_data=json.dumps(user_data),
            **profile_data,
        )
    except ValueError as e:
        return render_template("index.html", message="Invalid or expired token"), 401


from langchain_community.document_loaders import YoutubeLoader


@app.route("/get_transcripts", methods=["POST"])
def get_transcripts():
    data = request.get_json()
    youtube_url = data.get("youtube_url", "").strip()

    if not youtube_url:
        return jsonify({"error": "No YouTube URL provided"}), 400

    try:
        # Initialize the loader from the YouTube URL
        loader = YoutubeLoader.from_youtube_url(youtube_url, add_video_info=False)
        # Load transcript documents; each Document object contains a chunk of transcript text.
        transcript_docs = loader.load()

        # Combine the content from each Document into a single transcript string.
        transcript_text = " ".join(
            [doc.page_content for doc in transcript_docs]
        ).strip()

        if not transcript_text:
            logging.error("Transcript not found.")
            return jsonify({"error": "Transcript not found"}), 404

        return jsonify({"transcript": transcript_text})

    except Exception as e:
        logging.exception("Transcript extraction failed:")
        return jsonify({"error": f"Transcript extraction failed: {str(e)}"}), 500


@app.route("/scrape", methods=["POST"])
def scrape():
    data = request.get_json()
    url = data.get("url")
    if not url:
        return jsonify({"error": "Missing URL"}), 400
    try:
        response = requests.get(url)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        content = soup.get_text(separator=" ", strip=True)
        return jsonify({"success": True, "content": content})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/ai_suggs", methods=["GET", "POST"])
def chat_with_ai():
    email = session.get("email")
    user_data = get_user_data(email)
    if not user_data:
        return jsonify({"error": "User not found"}), 404
    print("yes")

    # Build the user's quiz history from their quiz_attempts
    quiz_attempts = user_data.get("quiz_attempts", [])
    quizzes_scores = []
    for attempt in quiz_attempts:
        topic = attempt.get("topic", "Unknown Topic")
        score = attempt.get("score", "Not Attempted")
        quizzes_scores.append({"topic": topic, "score": score})

    history = {"message": "here are your quizes and score:", "quizzes": quizzes_scores}
    # Print the result to the terminal for debugging
    print(json.dumps(history, indent=2))

    # Depending on the method, obtain the user question
    if request.method == "GET":
        user_question = request.args.get(
            "question", "What can you tell me about my quiz history?"
        )
    else:  # POST
        data = request.get_json()
        user_question = data.get(
            "message", "What can you tell me about my quiz history?"
        )

    # Build a prompt including the history and the user's question
    prompt = f"""
You are an advanced student assistant with access to the following user attempted quiz history. All scores are out of 10:
{json.dumps(history, indent=2)}

The user has asked the following question:
{user_question}

When answering, please:
1. Maintain a professional, precise, and helpful tone.
2. Refer to the user's quiz history where relevant.
3. After answering to the query, you may suggest to use features on this website such as uploading PDFs, dropping a website URL, embedding YouTube links, specifying custom topics, or using note descriptions from which quizzes can be generated to perpare for anything.
4. Be concise but thorough, focusing on actionable insights and recommendations.

Now, provide a short, professional response of not more than 100 words addressing the user’s question and any relevant app feature recommendations:
"""

    try:
        completion = client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert who guides students about their career based on their quiz results.",
                },
                {"role": "user", "content": prompt},
            ],
            temperature=0.3,
            max_tokens=1024,
            top_p=1,
            stream=False,
        )
        ai_answer = completion.choices[0].message.content.strip()
    except Exception as e:
        return (
            jsonify(
                {"error": f"An error occurred while processing AI suggestion: {e}"}
            ),
            500,
        )

    # For POST requests, return a simple JSON response for the chatbot
    if request.method == "POST":
        return jsonify({"response": ai_answer})

    # For GET requests, return the full details
    return jsonify(
        {
            "quiz_history": history,
            "user_question": user_question,
            "ai_answer": ai_answer,
        }
    )


from pptx import Presentation
from io import BytesIO


@app.route("/generate_ppt", methods=["POST"])
def generate_ppt():
    # 1) Parse payload
    data = request.get_json() or {}
    slides_count = int(data.get("slides", 1))
    topic = data.get("topic", "Presentation")
    description = data.get("description", "")

    # 2) LLaMA prompt
    llama_prompt = (
        f"Produce a JSON array of {slides_count} slides for a PowerPoint presentation on:\n"
        f"Topic: {topic}\nDescription: {description}\n\n"
        "Each array element must be an object with exactly two keys:\n"
        "  - title: the slide title\n"
        "  - content: an array of bullet strings OR a paragraph string with knowledgable and complete content about topic.\n\n"
        "Return ONLY valid JSON. Do NOT wrap it in markdown or any extra text."
    )

    try:
        # 3) Call LLaMA
        resp = client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {"role": "system", "content": "You draft PPT slide content as JSON."},
                {"role": "user", "content": llama_prompt},
            ],
            temperature=0.5,
            max_tokens=2048,
            top_p=1.0,
            stream=False,
        )
        raw = resp.choices[0].message.content.strip()
        slides_data = json.loads(raw)

        # 4) Build PPT
        prs = Presentation()
        # Title slide
        t0 = prs.slide_layouts[0]
        slide = prs.slides.add_slide(t0)
        slide.shapes.title.text = topic
        slide.placeholders[1].text = description

        # Content slides
        for idx, info in enumerate(slides_data[:slides_count], start=1):
            layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = info.get("title", f"Slide {idx+1}")
            tf = s.placeholders[1].text_frame
            # handle content as list or string
            content = info.get("content", "")
            if isinstance(content, list):
                lines = content
            else:
                lines = str(content).split("\n")
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                p = tf.add_paragraph()
                p.text = line
                p.level = 0

        # 5) Return PPTX
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(
            buf,
            as_attachment=True,
            download_name="presentation.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    except json.JSONDecodeError:
        logging.exception("LLaMA returned invalid JSON")
        return jsonify({"error": "Invalid JSON", "raw": raw}), 500
    except Exception as e:
        logging.exception("Error generating PPT")
        return jsonify({"error": str(e)}), 500


from flask import jsonify, request


@app.route("/sendQuizEmail", methods=["POST"])
@login_required  # optional — if you want only logged‑in users to hit it
def send_quiz_email():
    data = request.get_json() or {}
    email = data.get("email")
    quiz = data.get("quiz")
    print(email, quiz)
    if not email or not quiz:
        return jsonify(success=False, error="Missing payload"), 400

    try:
        msg = Message(
            subject="Your Quiz from QuizXpert",
            recipients=[email],
            body=f"Hello,\n\nHere’s your quiz:\n\n{quiz}\n\nGood luck!",
        )
        mail.send(msg)
        return jsonify(success=True)
    except Exception as e:
        app.logger.error(f"Error sending quiz email: {e}")
        return jsonify(success=False, error=str(e)), 500


from PIL import Image

import re
from google.api_core.exceptions import ResourceExhausted


@app.route("/get_quiz", methods=["POST"])
def get_quiz():
    if "quizPdf" not in request.files:
        return jsonify({"error": "No quizPdf file provided"}), 400

    file = request.files["quizPdf"]
    try:
        pdf_bytes = file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")

        quizzes = []
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            pix = page.get_pixmap()
            pil_img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)

            prompt = (
                "Extract all quiz questions and answers from this page image. "
                "Return ONLY a JSON array. Do not wrap it in markdown or code fences."
            )

            # ← Pass [prompt, pil_img] positionally
            try:
                response = gemini_model.generate_content([prompt, pil_img])
            except ResourceExhausted:
                return (
                    jsonify(
                        {
                            "error": "Gemini API quota exceeded—please wait or upgrade your plan."
                        }
                    ),
                    429,
                )
            except Exception as e:
                logging.exception(f"Gemini failed on page {page_num+1}:")
                return (
                    jsonify({"error": f"AI service error on page {page_num+1}: {e}"}),
                    500,
                )

            content = response.text

            # Strip any ```json fences
            cleaned = re.sub(r"^```(?:json)?\s*", "", content)
            cleaned = re.sub(r"\s*```$", "", cleaned).strip()

            try:
                page_quiz = json.loads(cleaned)
            except json.JSONDecodeError:
                logging.error(f"Failed to parse JSON on page {page_num+1}: {cleaned}")
                page_quiz = {"raw": cleaned}

            quizzes.append({"page": page_num + 1, "quiz": page_quiz})
            print(quizzes)

        return jsonify({"quizzes": quizzes}), 200

    except Exception as e:
        logging.exception("Error extracting quiz from PDF:")
        return jsonify({"error": f"Quiz extraction failed: {e}"}), 500


import fitz  # PyMuPDF
from flask import request, jsonify


@app.route("/validation", methods=["POST"])
def validation():
    # 1) Transcript (JSON preferred)
    data = request.get_json(silent=True) or {}
    transcript = (data.get("transcript") or request.form.get("transcript", "")).strip()
    if not transcript:
        return jsonify({"error": "Transcript is required"}), 400

    # 2) PDF → PIL.Image pages
    if "quizPdf" not in request.files:
        return jsonify({"error": "Quiz PDF is required"}), 400
    pdf_bytes = request.files["quizPdf"].read()
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    except Exception as e:
        return jsonify({"error": f"Invalid PDF: {e}"}), 400

    page_images = []
    for page in doc:
        pix = page.get_pixmap(dpi=150)
        png = pix.tobytes(output="png")
        pil_img = Image.open(io.BytesIO(png))
        page_images.append(pil_img)

    # 3) Build prompt
    prompt = f"You are a quiz validator, you have to check the user's quiz from the source i am providing, Respond with a JSON array of objects with keys: question, given_answer, is_correct, source including the exact location of answer. Selected answer may be right or wrong, you have to check slected answer from the source. the quiz: {page_images}, the source: {transcript}."

    # 4) Call Gemini
    try:
        response = gemini_model.generate_content([prompt] + page_images)
        content = getattr(response, "result", None) or getattr(response, "text", "")
        print(content)
    except ResourceExhausted:
        return jsonify({"error": "Quota exceeded"}), 429
    except Exception as e:
        logging.exception("Gemini API call failed")
        return jsonify({"error": str(e)}), 500

    # 5) Strip fences & parse JSON
    cleaned = re.sub(r"^```(?:json)?\s*", "", content)
    cleaned = re.sub(r"\s*```$", "", cleaned).strip()
    try:
        result = json.loads(cleaned)
    except json.JSONDecodeError:
        logging.error(f"Failed to parse JSON: {cleaned}")
        return jsonify({"error": "Invalid JSON from AI", "raw": cleaned}), 502

    return jsonify({"validation": result}), 200


@app.route("/validation_go")
def validation_go():
    return render_template("validation.html")


if __name__ == "__main__":
    app.run(debug=True)
